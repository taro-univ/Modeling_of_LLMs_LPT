#!/usr/bin/env python3
"""元資料の8・9・13枚目だけを修正した独立PowerPointを生成する。"""

from __future__ import annotations

import csv
import os
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/tmp/matplotlib-slide-revision-20260825")

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DECK_ROOT = Path(__file__).resolve().parents[1]
DATA_CSV = DECK_ROOT / "data" / "pancake_accuracy_600.csv"
FIGURE = DECK_ROOT / "figures" / "pancake_accuracy_600_heatmap.png"
OUTPUT = DECK_ROOT / "output" / "slides_08_09_13_revision_2026-08-25.pptx"

W = 13.333
H = 7.5

NAVY = "132238"
BLUE = "2D5BFF"
TEAL = "008A88"
RED = "D1495B"
ORANGE = "E7822B"
GRAY = "637083"
LIGHT = "F4F6F8"
LINE = "D9DEE5"
WHITE = "FFFFFF"
PALE_BLUE = "E8EDFF"
PALE_TEAL = "EAF7F6"

FONT_JP = "Meiryo UI"
FONT_LATIN = "Inter"


def rgb(value: str) -> RGBColor:
    """16進カラーをPowerPoint用RGBへ変換する。"""
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, fill=WHITE, line_color=None, radius=False):
    """矩形を追加する。"""
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line_color:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def line(slide, x1, y1, x2, y2, color=LINE, width=1.0, arrow=False):
    """直線を追加する。"""
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    if arrow:
        shape.line.end_arrowhead = True
    return shape


def text(
    slide, value, x, y, w, h, size=18, color=NAVY, bold=False,
    font=FONT_JP, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
    line_spacing=1.08,
):
    """文字列を追加する。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = valign
    for index, row in enumerate(str(value).split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = row
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = line_spacing
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return box


def bullet_list(slide, items, x, y, w, h, size=14, gap=7, accent=BLUE):
    """体言止め中心の箇条書きを追加する。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = "●  "
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.1
        paragraph.runs[0].font.name = FONT_JP
        paragraph.runs[0].font.size = Pt(size - 4)
        paragraph.runs[0].font.color.rgb = rgb(accent)
        run = paragraph.add_run()
        run.text = item
        run.font.name = FONT_JP
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(NAVY)
    return box


def base_slide(prs, title_value, original_page, title_size=25.5):
    """元資料と同じヘッダーを持つスライドを作る。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, WHITE)
    rect(slide, 0, 0, 0.16, H, BLUE)
    text(slide, title_value, 0.72, 0.55, 11.7, 0.5, title_size, NAVY, True)
    line(slide, 0.73, 1.146, 12.66, 1.146, LINE, 0.8)
    text(slide, "LLM REASONING DYNAMICS", 0.73, 7.14, 3.3, 0.18, 8.5, GRAY, True, FONT_LATIN)
    text(slide, f"{original_page:02d}", 11.9, 7.14, 0.75, 0.18, 8.5, GRAY, False, FONT_LATIN, PP_ALIGN.RIGHT)
    return slide


def source(slide, value):
    """出典をスライド下部へ追加する。"""
    text(slide, value, 0.78, 6.79, 11.8, 0.18, 7.7, GRAY, False, FONT_LATIN)


def load_accuracy_rows() -> list[dict[str, str]]:
    """Drive由来の600試行集計CSVを読み込む。"""
    with DATA_CSV.open(encoding="utf-8", newline="") as handle:
        rows = list(csv.DictReader(handle))
    if len(rows) != 6:
        raise ValueError(f"expected six cells, got {len(rows)}")
    if sum(int(row["n_trials"]) for row in rows) != 600:
        raise ValueError("trial total must be 600")
    for row in rows:
        expected = int(row["success_final"]) / int(row["n_trials"])
        if abs(float(row["accuracy"]) - expected) > 1e-12:
            raise ValueError(f"accuracy mismatch: {row['cell_id']}")
    return rows


def build_accuracy_figure(rows: list[dict[str, str]]) -> None:
    """N×min_movesのaccuracy heatmapを作成する。"""
    n_values = [3, 4, 5]
    mm_values = [3, 4, 5]
    matrix = np.full((3, 3), np.nan)
    counts: dict[tuple[int, int], tuple[int, int]] = {}
    for row in rows:
        n_value = int(row["N"])
        mm_value = int(row["min_moves"])
        matrix[n_values.index(n_value), mm_values.index(mm_value)] = float(row["accuracy"])
        counts[(n_value, mm_value)] = (int(row["success_final"]), int(row["n_trials"]))

    cmap = plt.colormaps["Blues"].copy()
    cmap.set_bad("#E5E7EB")
    figure, axis = plt.subplots(figsize=(7.0, 4.7), dpi=220)
    image = axis.imshow(matrix, vmin=0, vmax=1, cmap=cmap, aspect="equal")
    axis.set_xticks(range(3), [str(value) for value in mm_values], fontsize=12)
    axis.set_yticks(range(3), [str(value) for value in n_values], fontsize=12)
    axis.set_xlabel("min_moves", fontsize=13, fontweight="bold")
    axis.set_ylabel("N", fontsize=13, fontweight="bold", rotation=0, labelpad=18)
    axis.set_xticks(np.arange(-0.5, 3, 1), minor=True)
    axis.set_yticks(np.arange(-0.5, 3, 1), minor=True)
    axis.grid(which="minor", color="white", linewidth=3)
    axis.tick_params(which="minor", bottom=False, left=False)
    for row_index, n_value in enumerate(n_values):
        for col_index, mm_value in enumerate(mm_values):
            if np.isnan(matrix[row_index, col_index]):
                label = "N/A"
                color = "#6B7280"
            else:
                success, total = counts[(n_value, mm_value)]
                label = f"{matrix[row_index, col_index]:.0%}\n({success}/{total})"
                color = "white" if matrix[row_index, col_index] >= 0.65 else "#132238"
            axis.text(col_index, row_index, label, ha="center", va="center", fontsize=12, fontweight="bold", color=color)
    colorbar = figure.colorbar(image, ax=axis, fraction=0.048, pad=0.04)
    colorbar.set_label("final accuracy", fontsize=11)
    colorbar.set_ticks([0, 0.5, 1.0], labels=["0%", "50%", "100%"])
    for spine in axis.spines.values():
        spine.set_visible(False)
    figure.tight_layout()
    FIGURE.parent.mkdir(parents=True, exist_ok=True)
    figure.savefig(FIGURE, bbox_inches="tight", facecolor="white")
    plt.close(figure)


def add_hhmm_slide(prs):
    """修正版の元8枚目を追加する。"""
    slide = base_slide(
        prs,
        "2. 先行研究の紹介：Hidden Markov Modeling of Reasoning Dynamics",
        8,
        23.0,
    )
    text(slide, "Anonymous authors（ICLR 2026 submission）", 0.82, 1.37, 4.8, 0.28, 13.5, BLUE, True, FONT_LATIN)

    text(slide, "解析枠組み", 0.82, 1.82, 1.5, 0.3, 15, BLUE, True)
    rect(slide, 0.82, 2.22, 6.55, 3.75, LIGHT, None, True)
    text(slide, "explicit：semantic roleの遷移", 1.15, 2.5, 3.0, 0.3, 14, BLUE, True)
    roles = [("Setup", 1.16), ("Analysis", 2.52), ("Verify", 3.88), ("Final", 5.24)]
    for index, (name, xx) in enumerate(roles):
        rect(slide, xx, 3.02, 1.05, 0.58, WHITE, BLUE if name != "Verify" else ORANGE, True)
        text(slide, name, xx, 3.19, 1.05, 0.22, 11.5, NAVY, True, FONT_LATIN, PP_ALIGN.CENTER)
        if index < len(roles) - 1:
            line(slide, xx + 1.08, 3.31, roles[index + 1][1] - 0.05, 3.31, BLUE, 1.5, True)
    line(slide, 4.06, 3.82, 3.62, 3.82, RED, 1.3, True)
    text(slide, "verification loop", 3.28, 3.9, 1.6, 0.24, 10.5, RED, True, FONT_LATIN, PP_ALIGN.CENTER)

    text(slide, "implicit：layer-depth regime", 1.15, 4.52, 3.0, 0.3, 14, TEAL, True)
    depths = [("LOW", 1.22), ("MID", 2.65), ("UPPER", 4.08), ("ANCHOR", 5.51)]
    for index, (name, xx) in enumerate(depths):
        rect(slide, xx, 5.0, 1.02, 0.52, PALE_TEAL, TEAL, True)
        text(slide, name, xx, 5.15, 1.02, 0.2, 10.5, TEAL, True, FONT_LATIN, PP_ALIGN.CENTER)
        if index < len(depths) - 1:
            line(slide, xx + 1.05, 5.26, depths[index + 1][1] - 0.05, 5.26, TEAL, 1.4, True)

    text(slide, "手法", 7.73, 1.82, 1.0, 0.3, 15, BLUE, True)
    bullet_list(slide, [
        "生成文のreasoning step分割",
        "各step先頭tokenのlayer-wise hidden state",
        "PCA 64次元＋階層Hidden Markov Model",
        "semantic roleとdepth regimeの同時記述",
    ], 7.73, 2.22, 4.7, 1.78, 13.5, 6)
    line(slide, 7.73, 4.15, 12.15, 4.15, LINE, 0.7)
    text(slide, "主要結果", 7.73, 4.42, 1.3, 0.3, 15, TEAL, True)
    bullet_list(slide, [
        "推論早期のthink-first / commit-early分岐",
        "失敗軌道のverification loop・不安定なdepth遷移",
        "成功軌道の安定したsemantic path・構造的anchor",
    ], 7.73, 4.82, 4.7, 1.42, 13.5, 6, TEAL)
    source(slide, "Hidden Markov Modeling of Reasoning Dynamics in Large Language Models, OpenReview: fr9t7r43am")


def add_hypothesis_slide(prs):
    """8枚目の差し替えに合わせた元9枚目を追加する。"""
    slide = base_slide(prs, "3. 仮説：hidden-state dynamicsによる推論崩壊の記述", 9)
    rect(slide, 0.92, 1.48, 11.45, 1.62, PALE_BLUE, None, True)
    text(slide, "LLMの推論崩壊は、最終誤答より前のhidden-state trajectoryにおける", 1.28, 1.82, 10.72, 0.35, 18, NAVY, True, align=PP_ALIGN.CENTER)
    text(slide, "drift・滞留・final commit失敗として記述可能", 1.28, 2.38, 10.72, 0.4, 23, BLUE, True, align=PP_ALIGN.CENTER)
    links = [
        ("Shojaee", "制御された複雑度での崩壊", BLUE),
        ("Zhu", "誤り発生点の時間局在", TEAL),
        ("HHMM", "成功・失敗軌道の早期分岐", ORANGE),
    ]
    for index, (name, note, color) in enumerate(links):
        xx = 0.92 + index * 3.82
        rect(slide, xx, 3.65, 3.42, 1.18, LIGHT, None, True)
        text(slide, name, xx, 3.92, 3.42, 0.3, 15, color, True, FONT_LATIN, PP_ALIGN.CENTER)
        text(slide, note, xx + 0.18, 4.35, 3.06, 0.3, 13.5, NAVY, True, align=PP_ALIGN.CENTER)
    line(slide, 6.64, 4.95, 6.64, 5.34, BLUE, 1.7, True)
    rect(slide, 2.15, 5.38, 8.98, 0.72, PALE_TEAL, None, True)
    text(slide, "本発表の射程：挙動のモデリング　｜　制御モデル：将来課題", 2.15, 5.59, 8.98, 0.3, 17, TEAL, True, align=PP_ALIGN.CENTER)


def add_accuracy_slide(prs):
    """Drive上の600試行accuracyを示す元13枚目を追加する。"""
    slide = base_slide(prs, "4. 現状の共有：Nとmin_movesによる難易度分解", 13)
    text(slide, "T=0.6 / DeepSeek-R1-Distill-Qwen-14B / 600 trials（各cell n=100）", 0.82, 1.32, 7.3, 0.25, 11, GRAY, True, FONT_LATIN)
    slide.shapes.add_picture(str(FIGURE), Inches(0.72), Inches(1.62), width=Inches(6.3), height=Inches(4.77))

    text(slide, "min_moves=3固定", 7.42, 1.7, 3.5, 0.3, 15.5, BLUE, True, FONT_LATIN)
    text(slide, "N=3 → 4 → 5：100% → 86% → 76%", 7.42, 2.12, 4.75, 0.34, 17.5, NAVY, True, FONT_LATIN)
    text(slide, "同一計画長でもN増加に伴うaccuracy低下", 7.42, 2.58, 4.75, 0.32, 14, GRAY, True)
    line(slide, 7.42, 3.06, 12.12, 3.06, LINE, 0.8)

    text(slide, "N=5固定", 7.42, 3.38, 3.5, 0.3, 15.5, TEAL, True, FONT_LATIN)
    text(slide, "mm=3 → 4 → 5：76% → 38% → 38%", 7.42, 3.8, 4.75, 0.34, 17.5, NAVY, True, FONT_LATIN)
    text(slide, "mm=4での半減とmm=4–5間の横ばい", 7.42, 4.26, 4.75, 0.32, 14, GRAY, True)
    line(slide, 7.42, 4.74, 12.12, 4.74, LINE, 0.8)

    rect(slide, 7.42, 5.08, 4.7, 1.05, PALE_BLUE, None, True)
    text(slide, "層化によるN効果と計画長効果の分離", 7.7, 5.32, 4.14, 0.3, 16.5, BLUE, True, align=PP_ALIGN.CENTER)
    text(slide, "N=4ではmm4 > mm3：単純な単調則ではない結果", 7.7, 5.73, 4.14, 0.25, 12.5, NAVY, True, align=PP_ALIGN.CENTER)
    source(slide, "Source: Google Drive full_hidden_distribution_v1 / labels_v1.json, retrieved 2026-08-25; aggregate: data/pancake_accuracy_600.csv")


def make_deck() -> None:
    """独立した3枚の修正版PowerPointを生成する。"""
    rows = load_accuracy_rows()
    build_accuracy_figure(rows)
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "修正版スライド 08・09・13"
    prs.core_properties.subject = "元資料を変更しない独立差し替え用スライド"
    prs.core_properties.author = "LLM Reasoning Dynamics project"
    add_hhmm_slide(prs)
    add_hypothesis_slide(prs)
    add_accuracy_slide(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"saved: {OUTPUT}")
    print(f"slides: {len(prs.slides)} (original pages: 08, 09, 13)")


if __name__ == "__main__":
    make_deck()
