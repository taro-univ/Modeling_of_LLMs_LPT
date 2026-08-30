#!/usr/bin/env python3
"""2026-08-25の社内発表用PowerPointを生成する。

ユーザーが資料内へ記入した修正指示を反映した14枚構成。目次の文言を
各章タイトルの基準とし、追加情報はコロン以降のサブタイトルとして示す。
論文由来の主張には原図または公式ページの画像を添える。
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[4]
DECK_ROOT = Path(__file__).resolve().parents[1]
OUTPUT = DECK_ROOT / "output" / "hanoi_to_hidden_dynamics_2026-08-24.pptx"

SHOJAEE_FIG6 = DECK_ROOT / "figures" / "shojaee_2025_figure6.png"
SHOJAEE_FIG7 = DECK_ROOT / "figures" / "shojaee_2025_figure7.png"
ZHU_FIG2 = DECK_ROOT / "figures" / "zhu_2026_figure2.png"
CARSON_PAGE = DECK_ROOT / "figures" / "carson_2025_icml_page.png"
HANOI_FIG = DECK_ROOT / "figures" / "hanoi_deepseek_14b_accuracy.png"

W = 13.333
H = 7.5

NAVY = "132238"
BLUE = "2D5BFF"
TEAL = "008A88"
RED = "D1495B"
ORANGE = "E7822B"
GRAY = "637083"
LIGHT = "F4F6F8"
LINE = "D9DEE5"
WHITE = "FFFFFF"

FONT_JP = "Meiryo UI"
FONT_LATIN = "Inter"


def rgb(value: str) -> RGBColor:
    """16進カラーをPowerPoint用RGBへ変換する。"""
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, fill=WHITE, line_color=None, radius=False):
    """矩形を追加する。"""
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line_color:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def line(slide, x1, y1, x2, y2, color=LINE, width=1.0, arrow=False):
    """直線を追加する。"""
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    if arrow:
        shape.line.end_arrowhead = True
    return shape


def text(
    slide, value, x, y, w, h, size=18, color=NAVY, bold=False,
    font=FONT_JP, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
    margin=0.0, line_spacing=1.08,
):
    """文字列を追加する。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    for index, row in enumerate(str(value).split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = row
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = line_spacing
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return box


def bullet_list(slide, items, x, y, w, h, size=16, gap=8, accent=BLUE):
    """体言止め中心の箇条書きを追加する。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = "●  "
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.12
        paragraph.runs[0].font.name = FONT_JP
        paragraph.runs[0].font.size = Pt(size - 4)
        paragraph.runs[0].font.color.rgb = rgb(accent)
        run = paragraph.add_run()
        run.text = item
        run.font.name = FONT_JP
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(NAVY)
    return box


def add_image_contain(slide, path: Path, x, y, w, h):
    """画像を縦横比維持で指定領域へ収める。"""
    with Image.open(path) as image:
        iw, ih = image.size
    image_ratio = iw / ih
    frame_ratio = w / h
    if image_ratio >= frame_ratio:
        actual_w = w
        actual_h = w / image_ratio
        actual_x = x
        actual_y = y + (h - actual_h) / 2
    else:
        actual_h = h
        actual_w = h * image_ratio
        actual_x = x + (w - actual_w) / 2
        actual_y = y
    return slide.shapes.add_picture(
        str(path), Inches(actual_x), Inches(actual_y),
        width=Inches(actual_w), height=Inches(actual_h),
    )


def base_slide(prs, title_value, page):
    """2・3枚目に合わせた共通ヘッダーを持つスライドを作る。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, WHITE)
    rect(slide, 0, 0, 0.16, H, BLUE)
    text(slide, title_value, 0.72, 0.55, 11.7, 0.48, 25.5, NAVY, True)
    line(slide, 0.73, 1.146, 12.66, 1.146, LINE, 0.8)
    text(slide, "LLM REASONING DYNAMICS", 0.73, 7.14, 3.3, 0.18, 8.5, GRAY, True, FONT_LATIN)
    text(slide, f"{page:02d}", 11.9, 7.14, 0.75, 0.18, 8.5, GRAY, False, FONT_LATIN, PP_ALIGN.RIGHT)
    return slide


def source(slide, value, x=0.78, w=11.8):
    """図・主張の出典をスライド下部へ追加する。"""
    text(slide, value, x, 6.79, w, 0.18, 7.7, GRAY, False, FONT_LATIN)


def metric_row(slide, key, value, x, y, w, accent=BLUE):
    """実験条件または観測値を1行で表示する。"""
    text(slide, key, x, y, 1.55, 0.28, 11.5, accent, True, FONT_LATIN)
    text(slide, value, x + 1.6, y - 0.01, w - 1.6, 0.35, 14.5, NAVY, True)
    line(slide, x, y + 0.43, x + w, y + 0.43, LINE, 0.6)


def grouped_bar(slide, labels, final_values, search_values, x, y, w, h):
    """final accuracyとsearch goalの比較棒グラフを描く。"""
    plot_x = x + 0.65
    plot_y = y + 0.2
    plot_w = w - 0.9
    plot_h = h - 0.78
    line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, NAVY, 0.8)
    line(slide, plot_x, plot_y + plot_h, plot_x + plot_w, plot_y + plot_h, NAVY, 0.8)
    for tick in (0.0, 0.5, 1.0):
        yy = plot_y + plot_h * (1 - tick)
        line(slide, plot_x, yy, plot_x + plot_w, yy, LINE, 0.6)
        text(slide, f"{tick:.0%}", x, yy - 0.1, 0.52, 0.2, 8.5, GRAY, False, FONT_LATIN, PP_ALIGN.RIGHT)
    group_w = plot_w / len(labels)
    bar_w = 0.28
    for index, name in enumerate(labels):
        cx = plot_x + group_w * (index + 0.5)
        for offset, value, color in ((-0.17, final_values[index], BLUE), (0.17, search_values[index], TEAL)):
            bh = plot_h * value
            rect(slide, cx + offset - bar_w / 2, plot_y + plot_h - bh, bar_w, bh, color)
            text(slide, f"{value:.0%}", cx + offset - 0.32, plot_y + plot_h - bh - 0.25, 0.64, 0.2, 9, color, True, FONT_LATIN, PP_ALIGN.CENTER)
        text(slide, name, cx - 0.4, plot_y + plot_h + 0.1, 0.8, 0.23, 10.5, NAVY, True, FONT_LATIN, PP_ALIGN.CENTER)
    rect(slide, x + 1.0, y + h - 0.15, 0.16, 0.16, BLUE)
    text(slide, "final accuracy", x + 1.25, y + h - 0.19, 1.5, 0.22, 9.5, NAVY, False, FONT_LATIN)
    rect(slide, x + 2.85, y + h - 0.15, 0.16, 0.16, TEAL)
    text(slide, "search goal", x + 3.1, y + h - 0.19, 1.4, 0.22, 9.5, NAVY, False, FONT_LATIN)


def make_deck():
    """注記を反映した14枚の資料を生成する。"""
    for required in (SHOJAEE_FIG6, SHOJAEE_FIG7, ZHU_FIG2, CARSON_PAGE, HANOI_FIG):
        if not required.exists():
            raise FileNotFoundError(required)

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "大規模言語モデルの物理モデリング"
    prs.core_properties.subject = "2026-08-25 サマーインターン社内発表"
    prs.core_properties.author = "LLM Reasoning Dynamics project"

    # 1 表紙：ユーザーによる書き換えを保持
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, WHITE)
    rect(slide, 0, 0, 0.22, H, BLUE)
    text(slide, "SUMMER INTERNSHIP RESEARCH BRIEF", 0.85, 0.72, 5.8, 0.3, 12, BLUE, True, FONT_LATIN)
    text(slide, "大規模言語モデルの\n物理モデリング", 0.82, 1.62, 8.7, 1.35, 40, NAVY, True, line_spacing=1.0)
    text(slide, "推論崩壊をhidden-state dynamicsとして記述する試み", 0.86, 3.55, 8.6, 0.5, 19, GRAY)
    rect(slide, 9.86, 1.66, 2.4, 2.4, LIGHT, None, True)
    text(slide, "h(t)", 9.86, 2.08, 2.4, 0.55, 34, BLUE, True, FONT_LATIN, PP_ALIGN.CENTER)
    line(slide, 10.35, 3.02, 11.77, 3.02, TEAL, 3.0, True)
    text(slide, "drift / diffusion", 9.86, 3.34, 2.4, 0.3, 11, TEAL, True, FONT_LATIN, PP_ALIGN.CENTER)
    text(slide, "早稲田大学大学院　先進理工学研究科　物理学及応用物理学専攻", 0.86, 5.92, 8.4, 0.3, 12.5, NAVY, True)
    text(slide, "2026.08.25 | 社内発表", 0.86, 6.45, 4.5, 0.3, 12.5, NAVY, True)
    text(slide, "01", 11.85, 6.95, 0.5, 0.25, 9, GRAY, False, FONT_LATIN, PP_ALIGN.RIGHT)

    # 2 目次
    slide = base_slide(prs, "目次", 2)
    agenda = [
        ("01", "研究のモチベーション"),
        ("02", "先行研究の紹介"),
        ("03", "仮説"),
        ("04", "現状の共有"),
        ("05", "展望"),
    ]
    for index, (num, name) in enumerate(agenda):
        yy = 1.58 + index * 0.92
        text(slide, num, 1.02, yy, 0.7, 0.35, 17, BLUE if index < 3 else TEAL, True, FONT_LATIN)
        text(slide, name, 2.05, yy - 0.03, 6.0, 0.42, 21, NAVY, True)
        line(slide, 1.02, yy + 0.52, 12.0, yy + 0.52, LINE, 0.65)

    # 3 モチベーション1：Appleの観測事実
    slide = base_slide(prs, "1. 研究のモチベーション：複雑度に伴う推論崩壊", 3)
    add_image_contain(slide, SHOJAEE_FIG6, 0.72, 1.43, 8.7, 5.1)
    text(slide, "観測", 9.75, 1.62, 1.0, 0.3, 15, BLUE, True)
    bullet_list(slide, [
        "制御可能な4種のパズル",
        "複雑度増加に伴う正答率の急落",
        "モデル・問題固有の崩壊点",
    ], 9.75, 2.08, 2.65, 2.2, 14.5, 12)
    rect(slide, 9.72, 4.72, 2.72, 1.15, LIGHT, None, True)
    text(slide, "推論tokenの増加だけでは\n回避できない性能崩壊", 9.95, 5.0, 2.28, 0.62, 14.5, NAVY, True, align=PP_ALIGN.CENTER)
    source(slide, "Shojaee et al. (2025), The Illusion of Thinking, Fig. 6, arXiv:2506.06941v3, CC BY 4.0")

    # 4 モチベーション2：正解・不正解で異なる推論位置
    slide = base_slide(prs, "1. 研究のモチベーション：推論量の反転", 4)
    add_image_contain(slide, SHOJAEE_FIG7, 0.72, 1.52, 8.65, 4.65)
    text(slide, "同一複雑度内の差", 9.68, 1.64, 2.6, 0.3, 15, BLUE, True)
    bullet_list(slide, [
        "正解時の思考位置の後半化",
        "不正解時の思考位置の前半固定",
        "高複雑度での計算努力の減少",
    ], 9.68, 2.1, 2.7, 2.2, 14.5, 12)
    rect(slide, 9.65, 4.77, 2.78, 1.08, "E8EDFF", None, True)
    text(slide, "崩壊点以降における\n推論の早期打ち切り", 9.88, 5.03, 2.32, 0.6, 14.5, BLUE, True, align=PP_ALIGN.CENTER)
    source(slide, "Shojaee et al. (2025), Fig. 7: normalized position in thinking by puzzle complexity")

    # 5 モチベーション3：本研究への接続
    slide = base_slide(prs, "1. 研究のモチベーション：出力から内部状態へ", 5)
    text(slide, "既存の観測", 0.95, 1.64, 2.3, 0.35, 17, BLUE, True)
    rect(slide, 0.95, 2.14, 4.6, 2.55, LIGHT, None, True)
    text(slide, "accuracy", 1.32, 2.52, 1.7, 0.35, 22, NAVY, True, FONT_LATIN)
    text(slide, "thinking tokens", 1.32, 3.15, 2.4, 0.35, 22, NAVY, True, FONT_LATIN)
    text(slide, "text trajectory", 1.32, 3.78, 2.3, 0.35, 22, NAVY, True, FONT_LATIN)
    line(slide, 5.8, 3.42, 7.0, 3.42, BLUE, 2.8, True)
    text(slide, "未解明", 5.9, 2.85, 1.0, 0.3, 13, RED, True, align=PP_ALIGN.CENTER)
    text(slide, "本研究の観測対象", 7.27, 1.64, 3.0, 0.35, 17, TEAL, True)
    rect(slide, 7.27, 2.14, 5.0, 2.55, "EAF7F6", None, True)
    text(slide, "token-level hidden-state trajectory", 7.68, 2.55, 4.2, 0.4, 20, NAVY, True, FONT_LATIN, PP_ALIGN.CENTER)
    text(slide, "drift　／　滞留　／　final commit", 7.68, 3.36, 4.2, 0.35, 17, TEAL, True, align=PP_ALIGN.CENTER)
    text(slide, "最終誤答より前の内部変化", 7.68, 4.02, 4.2, 0.35, 16, NAVY, True, align=PP_ALIGN.CENTER)
    rect(slide, 0.95, 5.38, 11.32, 0.82, "E8EDFF", None, True)
    text(slide, "出力上の推論崩壊を、内部状態の時間発展として記述する研究", 0.95, 5.61, 11.32, 0.35, 19, BLUE, True, align=PP_ALIGN.CENTER)

    # 6 先行研究1：Shojaee et al.
    slide = base_slide(prs, "2. 先行研究の紹介：The Illusion of Thinking", 6)
    add_image_contain(slide, SHOJAEE_FIG6, 0.72, 1.5, 6.75, 4.72)
    text(slide, "Shojaee et al. (2025)", 7.8, 1.56, 3.6, 0.3, 15, BLUE, True, FONT_LATIN)
    metric_row(slide, "TASK", "Hanoiほか4種の制御可能なパズル", 7.8, 2.08, 4.55)
    metric_row(slide, "MODEL", "LRMと標準LLMの比較", 7.8, 2.68, 4.55)
    metric_row(slide, "SAMPLE", "各条件25試行・最大64k tokens", 7.8, 3.28, 4.55)
    text(slide, "主要結果", 7.8, 4.05, 1.4, 0.3, 14.5, TEAL, True)
    bullet_list(slide, [
        "低・中・高複雑度の3領域",
        "高複雑度での正答率の崩壊",
        "崩壊と同時に生じるthinking token減少",
    ], 7.8, 4.45, 4.55, 1.62, 14, 7, TEAL)
    source(slide, "Shojaee et al. (2025), arXiv:2506.06941v3; experimental setup and Fig. 6")

    # 7 先行研究2：Zhu et al.
    slide = base_slide(prs, "2. 先行研究の紹介：Dissecting Failure Dynamics", 7)
    add_image_contain(slide, ZHU_FIG2, 0.78, 1.45, 5.1, 5.1)
    text(slide, "Zhu et al. (ACL 2026)", 6.25, 1.55, 3.6, 0.3, 15, BLUE, True, FONT_LATIN)
    metric_row(slide, "MODEL", "DeepSeek-R1-Distill-Qwen-1.5B", 6.25, 2.05, 5.85)
    metric_row(slide, "DATA", "AMC / AIMEの推論trajectory", 6.25, 2.62, 5.85)
    metric_row(slide, "ORACLE", "segment単位のvalidity判定", 6.25, 3.19, 5.85)
    text(slide, "主要結果", 6.25, 3.92, 1.4, 0.3, 14.5, TEAL, True)
    bullet_list(slide, [
        "failure onsetの85%以上が生成前半30%以内",
        "誤りsegmentが1個のみのtrajectory 43.5%",
        "局所的entropy spike後の整合的だが誤った推論",
    ], 6.25, 4.34, 5.9, 1.55, 14, 7, TEAL)
    source(slide, "Zhu et al. (2026), Dissecting Failure Dynamics in LLM Reasoning, Fig. 2, ACL Anthology")

    # 8 先行研究3：Carson
    slide = base_slide(prs, "2. 先行研究の紹介：A Statistical Physics of Language Model Reasoning", 8)
    add_image_contain(slide, CARSON_PAGE, 0.78, 1.45, 4.85, 5.15)
    text(slide, "Carson (ICML 2025 Workshop)", 5.98, 1.55, 4.2, 0.3, 15, BLUE, True, FONT_LATIN)
    metric_row(slide, "STATE", "文単位・最終層のhidden state", 5.98, 2.05, 6.12)
    metric_row(slide, "SCALE", "8モデル × 7 reasoning benchmarks", 5.98, 2.62, 6.12)
    metric_row(slide, "MODEL", "drift–diffusion＋latent regime switching", 5.98, 3.19, 6.12)
    text(slide, "主要結果", 5.98, 3.92, 1.4, 0.3, 14.5, TEAL, True)
    bullet_list(slide, [
        "rank-40 drift manifoldによる約50%の分散説明",
        "4種のlatent reasoning regime",
        "switching linear dynamical systemによる再現",
    ], 5.98, 4.34, 6.12, 1.55, 14, 7, TEAL)
    source(slide, "Carson (2025), official ICML page: icml.cc/virtual/2025/50932")

    # 9 仮説
    slide = base_slide(prs, "3. 仮説：hidden-state dynamicsによる推論崩壊の記述", 9)
    rect(slide, 0.92, 1.48, 11.45, 1.62, "E8EDFF", None, True)
    text(slide, "LLMの推論崩壊は、最終誤答より前のhidden-state trajectoryにおける", 1.28, 1.82, 10.72, 0.35, 18, NAVY, True, align=PP_ALIGN.CENTER)
    text(slide, "drift・滞留・final commit失敗として記述可能", 1.28, 2.38, 10.72, 0.4, 23, BLUE, True, align=PP_ALIGN.CENTER)
    links = [
        ("Shojaee", "制御された複雑度での崩壊", BLUE),
        ("Zhu", "誤り発生点の時間局在", TEAL),
        ("Carson", "hidden軌道の確率過程モデル", ORANGE),
    ]
    for index, (name, note, color) in enumerate(links):
        xx = 0.92 + index * 3.82
        rect(slide, xx, 3.65, 3.42, 1.18, LIGHT, None, True)
        text(slide, name, xx, 3.92, 3.42, 0.3, 15, color, True, FONT_LATIN, PP_ALIGN.CENTER)
        text(slide, note, xx + 0.18, 4.35, 3.06, 0.3, 13.5, NAVY, True, align=PP_ALIGN.CENTER)
    line(slide, 6.64, 4.95, 6.64, 5.34, BLUE, 1.7, True)
    rect(slide, 2.15, 5.38, 8.98, 0.72, "EAF7F6", None, True)
    text(slide, "本発表の射程：挙動のモデリング　｜　制御モデル：将来課題", 2.15, 5.59, 8.98, 0.3, 17, TEAL, True, align=PP_ALIGN.CENTER)

    # 10 Hanoi実験設計
    slide = base_slide(prs, "4. 現状の共有：Hanoi実験の設計", 10)
    text(slide, "目的", 0.95, 1.55, 1.0, 0.3, 15, BLUE, True)
    text(slide, "Apple報告の推論崩壊をローカルopen modelで再現", 1.95, 1.52, 9.85, 0.38, 20, NAVY, True)
    line(slide, 0.95, 2.08, 12.1, 2.08, LINE, 0.7)
    columns = [
        ("問題複雑度", "N = 2–6", "最短手数 2ᴺ−1", BLUE),
        ("生成条件", "T = 0.1–3.0", "25–30 trials / cell", TEAL),
        ("観測量", "accuracy", "move loop / no move", ORANGE),
    ]
    for index, (head, value, note, color) in enumerate(columns):
        xx = 0.95 + index * 3.83
        text(slide, head, xx, 2.48, 3.42, 0.3, 14.5, color, True, align=PP_ALIGN.CENTER)
        rect(slide, xx, 2.94, 3.42, 1.45, LIGHT, None, True)
        text(slide, value, xx, 3.23, 3.42, 0.4, 23, NAVY, True, FONT_LATIN, PP_ALIGN.CENTER)
        text(slide, note, xx, 3.83, 3.42, 0.3, 13.5, GRAY, True, FONT_LATIN, PP_ALIGN.CENTER)
    rect(slide, 0.95, 4.92, 11.08, 1.1, "E8EDFF", None, True)
    text(slide, "出力", 1.28, 5.18, 0.8, 0.3, 14, BLUE, True)
    text(slide, "N–T平面上の正答率と支配的な失敗モード", 2.25, 5.16, 8.95, 0.35, 18, NAVY, True)
    text(slide, "対象：DeepSeek-R1-Distill-Qwen-14Bほか計5モデル", 2.25, 5.61, 8.95, 0.28, 13.5, GRAY, True)

    # 11 Hanoi結果
    slide = base_slide(prs, "4. 現状の共有：HanoiのN–T相図", 11)
    add_image_contain(slide, HANOI_FIG, 0.7, 1.48, 8.95, 4.95)
    text(slide, "DeepSeek-R1-Distill-Qwen-14B", 9.82, 1.58, 2.55, 0.48, 13, BLUE, True, FONT_LATIN)
    bullet_list(slide, [
        "N=2：T≤0.9で正答率100%",
        "N=3：48%以下からT=1.5で0%",
        "N=4：低温域20–36%、T=1.5以降0%",
        "N≥5：全観測温度で0%",
    ], 9.82, 2.28, 2.62, 2.72, 13.5, 10)
    rect(slide, 9.78, 5.28, 2.7, 0.93, "E8EDFF", None, True)
    text(slide, "N増加とT増加による\n単純な崩壊系列", 9.98, 5.5, 2.3, 0.5, 14.5, BLUE, True, align=PP_ALIGN.CENTER)
    source(slide, "Repository figure: figures/hanoi_nt_collapse/hanoi_nt_collapse_deepseek-r1-distill-qwen-14b.{png,csv}")

    # 12 Pancake実験設計
    slide = base_slide(prs, "4. 現状の共有：Pancake Sortingの実験設計", 12)
    text(slide, "Hanoiの制約", 0.95, 1.52, 2.1, 0.32, 15, RED, True)
    rect(slide, 0.95, 1.98, 4.7, 1.38, LIGHT, None, True)
    text(slide, "N増加に伴う最短手数の指数増加", 1.25, 2.3, 4.1, 0.34, 18, NAVY, True, align=PP_ALIGN.CENTER)
    text(slide, "状態表現と計画長の交絡", 1.25, 2.82, 4.1, 0.3, 15, GRAY, True, align=PP_ALIGN.CENTER)
    line(slide, 5.85, 2.66, 6.95, 2.66, BLUE, 2.6, True)
    text(slide, "Pancakeの利点", 7.2, 1.52, 2.4, 0.32, 15, TEAL, True)
    rect(slide, 7.2, 1.98, 5.02, 1.38, "EAF7F6", None, True)
    text(slide, "Nと厳密なmin_movesの層化", 7.5, 2.3, 4.42, 0.34, 18, NAVY, True, align=PP_ALIGN.CENTER)
    text(slide, "表現長と計画長の独立評価", 7.5, 2.82, 4.42, 0.3, 15, TEAL, True, align=PP_ALIGN.CENTER)
    metric_row(slide, "MODEL", "DeepSeek-R1-Distill-Qwen-14B", 0.95, 4.02, 5.35)
    metric_row(slide, "TEMP", "T = 0.6", 0.95, 4.64, 5.35)
    metric_row(slide, "GRID", "N = 3–5 / min_moves = 3–5", 6.82, 4.02, 5.4, TEAL)
    metric_row(slide, "LABEL", "search goalとfinal answerの分離", 6.82, 4.64, 5.4, TEAL)
    rect(slide, 0.95, 5.55, 11.27, 0.63, "E8EDFF", None, True)
    text(slide, "18 instancesによるdebug sweep　→　outcome別hidden軌道の取得設計", 0.95, 5.73, 11.27, 0.3, 16.5, BLUE, True, align=PP_ALIGN.CENTER)

    # 13 Pancake結果：ユーザーが「結構いい」とした構成を保持
    slide = base_slide(prs, "4. 現状の共有：Nとmin_movesによる難易度分解", 13)
    text(slide, "T=0.6 / DeepSeek-R1-Distill-Qwen-14B / 18 instances", 0.82, 1.32, 7.0, 0.25, 11, GRAY, True, FONT_LATIN)
    grouped_bar(slide, ["N=3", "N=4", "N=5"], [1.00, 0.83, 0.44], [1.00, 0.83, 0.67], 0.72, 1.62, 6.2, 4.85)
    text(slide, "同じ min_moves=3", 7.45, 1.78, 3.6, 0.32, 16, BLUE, True, FONT_LATIN)
    text(slide, "N=3 → N=4/5での\nfinal accuracy低下", 7.45, 2.22, 4.4, 0.72, 18, NAVY, True)
    line(slide, 7.45, 3.12, 12.1, 3.12, LINE, 0.8)
    text(slide, "同じ N=5", 7.45, 3.46, 3.6, 0.32, 16, TEAL, True, FONT_LATIN)
    text(slide, "min_moves 3 → 4/5での\nfinal accuracy 67% → 33%", 7.45, 3.9, 4.4, 0.72, 18, NAVY, True)
    rect(slide, 7.45, 5.07, 4.65, 0.82, LIGHT, None, True)
    text(slide, "search goal > final accuracy", 7.75, 5.29, 3.95, 0.28, 16, RED, True, FONT_LATIN, PP_ALIGN.CENTER)
    text(slide, "探索成功後のfinal commit失敗", 7.45, 6.12, 4.65, 0.27, 13.5, NAVY, True, align=PP_ALIGN.CENTER)
    source(slide, "Source: docs/research_state/results_summary.md — min_moves stratified debug sweep")

    # 14 展望・まとめ
    slide = base_slide(prs, "5. 展望：挙動のモデリングから制御へ", 14)
    findings = [
        ("先行研究", "崩壊点・failure onset・確率過程という三つの記述", BLUE),
        ("Hanoi", "Apple報告の推論崩壊をローカルopen modelで再現", TEAL),
        ("Pancake", "Nとmin_movesの分離による挙動モデリングの土台", ORANGE),
    ]
    for index, (tag, value, color) in enumerate(findings):
        yy = 1.55 + index * 1.03
        text(slide, tag, 0.95, yy, 1.45, 0.32, 14.5, color, True)
        text(slide, value, 2.6, yy - 0.02, 9.4, 0.38, 18, NAVY, True)
        line(slide, 0.95, yy + 0.54, 12.05, yy + 0.54, LINE, 0.7)
    text(slide, "次段階", 0.95, 4.8, 1.45, 0.32, 15, BLUE, True)
    rect(slide, 2.58, 4.66, 9.47, 1.26, "E8EDFF", None, True)
    text(slide, "複数trajectoryからのoutcome別drift / diffusion推定", 2.92, 4.93, 8.8, 0.34, 18, BLUE, True, align=PP_ALIGN.CENTER)
    text(slide, "崩壊の早期予測　→　将来的なinference-time control", 2.92, 5.42, 8.8, 0.3, 15.5, NAVY, True, align=PP_ALIGN.CENTER)
    text(slide, "本発表の到達点：制御ではなく、推論崩壊の動力学的診断", 0.95, 6.35, 11.1, 0.32, 17, TEAL, True, align=PP_ALIGN.CENTER)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"saved: {OUTPUT}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    make_deck()
